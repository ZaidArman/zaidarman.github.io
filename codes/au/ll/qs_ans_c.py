import numpy as  np
import pandas as pd

from pptx import Presentation
import pptx
import inspect
import copy
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR


def ans(serial,i,dfx,ppt):
    TEMPLATE_SLIDE_QS = ppt.slides[0]
    TEMPLATE_SLIDE_ANS_MATH = ppt.slides[0]

    slide_layout_QS = ppt.slide_layouts[0]
    slide_layout_ANS = ppt.slide_layouts[0]
    slide_layout_ANS_MATH  = ppt.slide_layouts[0]

    NEW_SLIDE = ppt.slides.add_slide(slide_layout_ANS_MATH)
    
    for shp in TEMPLATE_SLIDE_ANS_MATH.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        NEW_SLIDE.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
    for shape in NEW_SLIDE.shapes: 
        if not shape.has_text_frame: 
            continue 
        for paragraph in shape.text_frame.paragraphs: 
            
            for run in paragraph.runs:
                if (dfx['Answer'][i]==1) and ('01' in run.text):
                    font = run.font
                    font.bold = True 
                    font.color.rgb =ans_color
                    shape.fill.fore_color.rgb = ans_box_color
                    
                elif (dfx['Answer'][i]==2) and ('02' in run.text):
                    font = run.font
                    font.bold = True 
                    font.color.rgb = ans_color
                    shape.fill.fore_color.rgb  = ans_box_color
                    
                elif (dfx['Answer'][i]==3) and ('03' in run.text):
                    font = run.font
                    font.bold = True 
                    font.color.rgb = ans_color
                    shape.fill.fore_color.rgb  = ans_box_color
                    
                elif (dfx['Answer'][i]==4) and ('04' in run.text):
                    font = run.font
                    font.bold = True 
                    font.color.rgb = ans_color
                    shape.fill.fore_color.rgb  = ans_box_color

                    
            for run in paragraph.runs: 
                if  'QuestionNo' in run.text :
                    run.text = "Question: "+str(serial)
                elif  'Question' in run.text :
                    run.text = ""+str(dfx['Question'][i]).strip()
                elif '01' in run.text :
                    run.text =  "ক) "+ str(dfx['Option A'][i]).strip()
                elif '02' in run.text :
                    run.text = "খ) "+str(dfx['Option B'][i]).strip()
                elif '03' in run.text :
                    run.text = "গ) "+str(dfx['Option C'][i]).strip()
                elif '04' in run.text :
                    run.text = "ঘ) "+str(dfx['Option D'][i]).strip()


                           
    NEW_SLIDE =False  
    return ppt



def content(serial,i,dfx,ppt):
    TEMPLATE_SLIDE_QS = ppt.slides[0]
    TEMPLATE_SLIDE_ANS_MATH = ppt.slides[0]

    slide_layout_QS = ppt.slide_layouts[0]
    slide_layout_ANS = ppt.slide_layouts[0]
    slide_layout_ANS_MATH  = ppt.slide_layouts[0]
    NEW_SLIDE = ppt.slides.add_slide(slide_layout_ANS_MATH)
    
    for shp in TEMPLATE_SLIDE_ANS_MATH.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        NEW_SLIDE.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
    for shape in NEW_SLIDE.shapes: 
        if not shape.has_text_frame: 
            continue 
        for paragraph in shape.text_frame.paragraphs: 
            for run in paragraph.runs: 
                if  'Title' in run.text :
                    run.text = ""+str(dfx['Title'][i]).strip()
                elif  'Content' in run.text :
                    run.text = ""+str(dfx['Content'][i]).strip()
    NEW_SLIDE =False  
    serial=serial+1

    
    return ppt,serial



def qs(serial,i,dfx,ppt):
    TEMPLATE_SLIDE_QS = ppt.slides[0]
    TEMPLATE_SLIDE_ANS_MATH = ppt.slides[0]

    slide_layout_QS = ppt.slide_layouts[0]
    slide_layout_ANS = ppt.slide_layouts[0]
    slide_layout_ANS_MATH  = ppt.slide_layouts[0]
    NEW_SLIDE = ppt.slides.add_slide(slide_layout_ANS_MATH)
    
    for shp in TEMPLATE_SLIDE_ANS_MATH.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        NEW_SLIDE.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
    for shape in NEW_SLIDE.shapes: 
        if not shape.has_text_frame: 
            continue 
        for paragraph in shape.text_frame.paragraphs: 
            for run in paragraph.runs: 
                if  'QuestionNo' in run.text :
                    run.text = "Question: "+str(serial)
                elif  'Question' in run.text :
                    run.text = ""+str(dfx['Question'][i]).strip()
                elif '01' in run.text :
                    run.text =  "A) "+ str(dfx['Option A'][i]).strip()
                elif '02'in run.text :
                    run.text = "B) "+str(dfx['Option B'][i]).strip()
                elif '03' in run.text :
                    run.text = "C) "+str(dfx['Option C'][i]).strip()
                elif '04'in run.text :
                    run.text = "D) "+str(dfx['Option D'][i]).strip()
    NEW_SLIDE =False  
    
    ans(serial,i,dfx,ppt)
    serial=serial+1

    
    return ppt,serial